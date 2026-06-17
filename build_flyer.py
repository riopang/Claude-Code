"""
Door-knock leave-behind flyer — Thomson Garden Estate edition
A4 portrait, navy + gold palette, with agent photo
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
GOLD    = RGBColor(0xC9, 0xA8, 0x4C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF4, 0xF6, 0xFA)
LTGOLD  = RGBColor(0xFD, 0xF6, 0xE3)
MIDGRAY = RGBColor(0x64, 0x74, 0x8B)
DARK2   = RGBColor(0x16, 0x23, 0x38)
LTBLUE  = RGBColor(0xB0, 0xBE, 0xD4)
RULE    = RGBColor(0xD9, 0xE2, 0xEF)

W = 8.27
H = 11.69

PHOTO_PATH = "/Users/riopang/Desktop/Claude Code/rio_face_square.png"

def add_rect(slide, x, y, w, h, fill, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def add_oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(1.5)
    else: s.line.fill.background()
    return s

def tb(slide, x, y, w, h, text, size=11, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = "Calibri"
    if color: r.font.color.rgb = color
    return box

def set_bg(slide, color):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = color

def build():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # ── HEADER BAND ─────────────────────────────────────────────────────────────
    HDR_H = 2.65
    add_rect(s, 0, 0, W, HDR_H, NAVY)
    add_rect(s, 0, 0, 0.28, HDR_H, GOLD)   # left gold bar

    # Headline
    tb(s, 0.45, 0.22, 5.8, 0.72,
       "Is now the right time\nto sell your landed home?",
       size=24, bold=True, color=WHITE)

    # Sub — estate-specific
    tb(s, 0.45, 1.02, 5.8, 0.4,
       "Thomson Garden Estate has seen 21 transactions in the last 12 months — "
       "and serious buyers are actively looking here.",
       size=10.5, italic=True, color=GOLD)

    # ── AGENT CARD (right side of header) ──────────────────────────────────────
    CARD_X, CARD_Y, CARD_W, CARD_H = 6.35, 0.20, 1.72, 2.40
    add_rect(s, CARD_X, CARD_Y, CARD_W, CARD_H, DARK2, RGBColor(0x2E, 0x41, 0x66), Pt(0.75))

    tb(s, CARD_X+0.05, CARD_Y+0.06, CARD_W-0.10, 0.20,
       "ERA SINGAPORE", size=8.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_rect(s, CARD_X+0.15, CARD_Y+0.28, CARD_W-0.30, 0.03, GOLD)

    # Agent photo (square, centered in card)
    PHOTO_W = 1.38
    photo_x = CARD_X + (CARD_W - PHOTO_W) / 2
    s.shapes.add_picture(PHOTO_PATH, Inches(photo_x), Inches(CARD_Y + 0.33),
                         width=Inches(PHOTO_W))

    # Name, CEA, phone below photo
    TEXT_Y = CARD_Y + 0.33 + PHOTO_W + 0.06   # photo is square so height ≈ width
    tb(s, CARD_X+0.05, TEXT_Y,      CARD_W-0.10, 0.22,
       "Rio Pang", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, CARD_X+0.05, TEXT_Y+0.22, CARD_W-0.10, 0.18,
       "CEA No. R0702148", size=8, color=LTBLUE, align=PP_ALIGN.CENTER)
    tb(s, CARD_X+0.05, TEXT_Y+0.40, CARD_W-0.10, 0.22,
       "9183 1316", size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ── 4 STAT BOXES ────────────────────────────────────────────────────────────
    stats = [
        ("+6.7%",  "landed price growth\nyear-on-year (2025)"),
        ("+73%",   "URA landed price index\nQ1 2019 → Q1 2026"),
        ("21 deals","Thomson Garden Estate\nlast 12 months"),
        ("+46%",   "ultra-wealthy buyer count\nforecast by 2031"),
    ]
    STATS_Y = HDR_H + 0.12
    bw = (W - 1.0) / 4
    for i, (val, lbl) in enumerate(stats):
        x = 0.5 + i * bw
        add_rect(s, x, STATS_Y, bw - 0.12, 1.12, LTGOLD, GOLD, Pt(1))
        tb(s, x+0.06, STATS_Y+0.07, bw-0.24, 0.48, val,
           size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tb(s, x+0.06, STATS_Y+0.54, bw-0.24, 0.5, lbl,
           size=8.5, color=MIDGRAY, align=PP_ALIGN.CENTER)

    # ── WHY NOW ─────────────────────────────────────────────────────────────────
    WHY_Y = STATS_Y + 1.25
    tb(s, 0.5, WHY_Y, 4.8, 0.28, "WHY SELLERS ARE MOVING NOW",
       size=8.5, bold=True, color=GOLD)
    add_rect(s, 0.5, WHY_Y+0.3, 4.8, 2.62, LIGHT, RULE, Pt(0.5))

    reasons = [
        ("Prices up 6.7% YoY",         "Strong market recovery — Q2/Q3 is peak transaction season."),
        ("Declining interest rates",    "SORA falling → buyer affordability improving → larger qualified buyer pool."),
        ("Wealthiest buyer base ever",  "Singapore UHNWI count up 55% in 5 years. Forecast: +46% more by 2031."),
        ("Supply stays scarce",         "Landed stock grew only 12% in 25 years vs 228% for non-landed condos."),
        ("60% ABSD for foreigners",     "Landed stays exclusive to locals/PRs — buyers compete in a motivated pool."),
    ]
    for i, (title, desc) in enumerate(reasons):
        y = WHY_Y + 0.40 + i * 0.48
        add_oval(s, 0.62, y+0.06, 0.26, 0.26, NAVY)
        tb(s, 0.62, y+0.06, 0.26, 0.26, str(i+1), size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        tb(s, 0.98, y+0.04, 1.4, 0.22, title, size=9.5, bold=True, color=NAVY)
        tb(s, 0.98, y+0.25, 3.2, 0.2,  desc,  size=8.5, color=MIDGRAY)

    # ── TGE PRICE DATA ──────────────────────────────────────────────────────────
    PRICE_X = 5.5
    tb(s, PRICE_X, WHY_Y, 2.58, 0.28, "THOMSON GARDEN ESTATE",
       size=8, bold=True, color=GOLD)
    tb(s, PRICE_X, WHY_Y+0.18, 2.58, 0.14, "Jun 2025 – Jun 2026 · Freehold D20",
       size=7, italic=True, color=MIDGRAY)
    add_rect(s, PRICE_X, WHY_Y+0.30, 2.58, 2.62, LIGHT, RULE, Pt(0.5))

    prices = [
        ("Terrace (typical)",  "$2.65M–$4.35M"),
        ("Terrace (large)",    "$4.35M–$6.25M"),
        ("Semi-D",             "$5.70M–$7.20M"),
        ("Avg Land PSF",       "$3,151 psf"),
    ]
    for i, (label, price) in enumerate(prices):
        y = WHY_Y + 0.40 + i * 0.60
        add_rect(s, PRICE_X+0.10, y, 2.36, 0.50, WHITE, RULE, Pt(0.5))
        tb(s, PRICE_X+0.18, y+0.05, 1.2, 0.24, label, size=9, color=NAVY)
        tb(s, PRICE_X+0.18, y+0.27, 2.0, 0.18, price,
           size=12, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    tb(s, PRICE_X, WHY_Y+3.0, 2.58, 0.18,
       "Source: URA REALIS (Jun 2025–Jun 2026)",
       size=7, italic=True, color=MIDGRAY)

    # ── WHAT YOUR HOME COULD BE WORTH ───────────────────────────────────────────
    WORTH_Y = WHY_Y + 2.98
    add_rect(s, 0, WORTH_Y, W, 1.32, NAVY)
    add_rect(s, 0, WORTH_Y, 0.28, 1.32, GOLD)
    tb(s, 0.45, WORTH_Y+0.08, 5.5, 0.28, "WHAT YOUR HOME COULD BE WORTH",
       size=9, bold=True, color=GOLD)
    tb(s, 0.45, WORTH_Y+0.38, 7.5, 0.76,
       "Based on 21 transactions in Thomson Garden Estate over the last 12 months, "
       "freehold terraces have transacted between $2.65M–$6.25M and semi-detached "
       "homes at $5.70M–$7.20M. Average land PSF: $3,151. "
       "Call me for a complimentary, no-obligation valuation of your specific home.",
       size=10, color=WHITE)

    # ── 3-STEP PROCESS ──────────────────────────────────────────────────────────
    STEPS_Y = WORTH_Y + 1.40
    tb(s, 0.5, STEPS_Y, 7.5, 0.28, "HOW IT WORKS — THREE SIMPLE STEPS",
       size=8.5, bold=True, color=GOLD)

    steps = [
        ("1", "Free property assessment",   "I visit, assess condition, pull the latest comps and give you an honest, no-obligation valuation."),
        ("2", "Marketing strategy meeting", "We agree on the right price and approach. I'll show you exactly who the buyers are and how I'll reach them."),
        ("3", "Sign listing & launch",      "Professional photography within 5 days. Listing live within 7 days. Active buyer outreach starts immediately."),
    ]
    sw = (W - 1.0) / 3
    for i, (num, title, desc) in enumerate(steps):
        x = 0.5 + i * sw
        add_rect(s, x, STEPS_Y+0.32, sw - 0.15, 1.55, LTGOLD, GOLD, Pt(0.75))
        add_oval(s, x + (sw-0.15)/2 - 0.20, STEPS_Y+0.39, 0.40, 0.40, NAVY)
        tb(s, x + (sw-0.15)/2 - 0.20, STEPS_Y+0.39, 0.40, 0.40, num,
           size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        tb(s, x+0.10, STEPS_Y+0.86, sw-0.35, 0.30, title,
           size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tb(s, x+0.10, STEPS_Y+1.16, sw-0.35, 0.65, desc,
           size=8.5, color=MIDGRAY, align=PP_ALIGN.CENTER)

    # ── FOOTER BAND ─────────────────────────────────────────────────────────────
    FOOTER_Y = 10.52
    add_rect(s, 0, FOOTER_Y, W, 1.08, DARK2)
    add_rect(s, 0, FOOTER_Y, 0.28, 1.08, GOLD)

    tb(s, 0.45, FOOTER_Y+0.08, 3.5, 0.28, "Rio Pang 彭仕轩",
       size=13, bold=True, color=WHITE)
    tb(s, 0.45, FOOTER_Y+0.38, 3.5, 0.24,
       "ERA Singapore  |  CEA Reg. No. R0702148",
       size=9, color=LTBLUE)

    tb(s, 4.0, FOOTER_Y+0.08, 2.5, 0.26, "9183 1316", size=13, bold=True, color=GOLD)
    tb(s, 4.0, FOOTER_Y+0.36, 2.5, 0.26, "WhatsApp / Call", size=9.5, color=LTBLUE)

    # Disclaimer
    tb(s, 0.45, FOOTER_Y+0.85, W-0.9, 0.20,
       "Data sourced from URA REALIS (Jun 2025–Jun 2026). This flyer is for informational purposes only and does not constitute a formal valuation.",
       size=6.5, italic=True, color=MIDGRAY)

    out = "/Users/riopang/Desktop/Claude Code/Doorknock Leave-Behind Flyer.pptx"
    prs.save(out)
    print(f"Saved → {out}")

build()
