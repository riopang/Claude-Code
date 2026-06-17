"""
Landed Seller Pitch Deck — Singapore 2026
16-slide PowerPoint, navy + gold palette, Calibri font
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import lxml.etree as etree

# ─── COLOURS ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
GOLD    = RGBColor(0xC9, 0xA8, 0x4C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF4, 0xF6, 0xFA)
LTGOLD  = RGBColor(0xFD, 0xF6, 0xE3)
MIDGRAY = RGBColor(0x64, 0x74, 0x8B)
PLHOLD  = RGBColor(0xE8, 0xED, 0xF5)
RULE    = RGBColor(0xD9, 0xE2, 0xEF)
DARK2   = RGBColor(0x16, 0x23, 0x38)
DARK3   = RGBColor(0x0F, 0x1D, 0x32)
LTBLUE  = RGBColor(0xB0, 0xBE, 0xD4)

# ─── HELPERS ──────────────────────────────────────────────────────────────────
def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, fill_color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        9,  # oval
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_size=12, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, font_face="Calibri",
                word_wrap=True, v_anchor=None):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    # clear default paragraph
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_face
    if color:
        run.font.color.rgb = color
    return txBox

def add_label_title(slide, label, title, dark=False):
    label_color = GOLD
    title_color = WHITE if dark else NAVY
    # Label
    add_textbox(slide, 0.5, 0.25, 9.0, 0.28, label.upper(),
                font_size=8.5, bold=True, color=label_color, font_face="Calibri")
    # Title
    add_textbox(slide, 0.5, 0.50, 9.0, 0.72, title,
                font_size=26, bold=True, color=title_color, font_face="Calibri")

def add_stat_box(slide, val, label, x, y, w=2.7, h=1.18):
    add_rect(slide, x, y, w, h, LTGOLD, GOLD, Pt(1.5))
    add_textbox(slide, x+0.12, y+0.05, w-0.24, 0.58, val,
                font_size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, x+0.12, y+0.62, w-0.24, 0.5, label,
                font_size=9.5, color=MIDGRAY, align=PP_ALIGN.CENTER)

def add_bullets(slide, items, x, y, w, h, font_size=12):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet via XML
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', 'Arial')
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts.set('val', '600')
        ind = etree.SubElement(pPr, qn('a:ind'))
        ind.set('marL', str(int(Inches(0.22))))
        ind.set('indent', str(int(Inches(-0.22))))
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.name = "Calibri"
        run.font.color.rgb = NAVY

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]  # blank

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 1 — COVER
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, NAVY)

    # Left gold accent bar
    add_rect(s, 0, 0, 0.3, 5.625, GOLD)

    # Headline
    add_textbox(s, 0.55, 0.85, 7.6, 1.9,
                "Thinking of Selling\nYour Landed Home?",
                font_size=38, bold=True, color=WHITE)

    # Sub-headline
    add_textbox(s, 0.55, 2.82, 7.2, 0.45,
                "Your trusted landed property specialist in Singapore.",
                font_size=15, italic=True, color=GOLD)

    # Divider
    add_rect(s, 0.55, 3.38, 4.4, 0.04, GOLD)

    # Tagline
    add_textbox(s, 0.55, 3.5, 7.0, 0.38,
                "Data-driven.  Results-focused.  Landed specialist.",
                font_size=12, color=LTBLUE)

    # Agent info box
    add_rect(s, 0.55, 4.05, 5.0, 1.25, DARK2, RGBColor(0x2E,0x41,0x66), Pt(1))
    add_textbox(s, 0.78, 4.15, 4.6, 0.38,
                "[INSERT: Agent Name]", font_size=14, bold=True, color=WHITE)
    add_textbox(s, 0.78, 4.52, 4.6, 0.28,
                "[INSERT: CEA Reg. No.]", font_size=10, color=LTBLUE)
    add_textbox(s, 0.78, 4.78, 4.6, 0.28,
                "[INSERT: Mobile]  |  [INSERT: Email]", font_size=10, color=GOLD)

    # ERA badge
    add_rect(s, 8.1, 4.05, 1.6, 1.25, DARK2, RGBColor(0x2E,0x41,0x66), Pt(1))
    add_textbox(s, 8.1, 4.35, 1.6, 0.65, "ERA\nSINGAPORE",
                font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 2 — SCARCITY
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Market Context", "Singapore Landed — A Structurally Scarce Asset")

    add_stat_box(s, "~75,000",  "landed homes in Singapore\n(~5–6% of housing stock)", 0.5,  1.45, 2.85, 1.3)
    add_stat_box(s, "+12%",     "supply growth over 25 years\n(vs. +228% for non-landed condos)", 3.57, 1.45, 2.85, 1.3)
    add_stat_box(s, "+73%",     "URA landed price index\nQ1 2019 → Q1 2026", 6.62, 1.45, 2.85, 1.3)

    add_bullets(s, [
        "Supply grew from 67,229 units (2000) to 75,338 units (2025) — just 12% over 25 years.",
        "Non-landed condos grew 228% over the same period: 114,532 → 375,612 units.",
        "The government releases very little new land for landed housing — this scarcity is permanent and structural.",
        "Buyers recognise this: landed homes are a defensive, inflation-resistant asset that supply cannot scale to meet.",
        "Result: landed prices have compounded even through multiple rounds of government cooling measures.",
    ], 0.5, 2.88, 9.0, 2.4)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 3 — MARKET SNAPSHOT
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Market Data", "2026 Landed Market Snapshot")

    add_stat_box(s, "+6.7%",  "YoY price growth\n(landed, full year 2025)",  0.5,  1.45, 2.2, 1.08)
    add_stat_box(s, "+3.4%",  "QoQ surge in Q4 2025\n(strongest quarter)",   2.82, 1.45, 2.2, 1.08)
    add_stat_box(s, "+15%",   "transaction value growth\n(2025 vs 2024 YoY)", 5.14, 1.45, 2.2, 1.08)
    add_stat_box(s, "+5–7%",  "ERA 2026 price forecast\n(1,750–1,950 units)", 7.46, 1.45, 2.2, 1.08)

    # Table
    rows_data = [
        ("Property Type",           "Region", "Median Price (2026)", True),
        ("Terrace (Intermediate)",  "OCR",    "$4.44M",              False),
        ("Semi-Detached",           "OCR",    "$5.75M",              False),
        ("Detached",                "RCR",    "$10.9M",              False),
        ("Detached / GCB",          "CCR",    "$16.825M",            False),
    ]
    tbl = s.shapes.add_table(5, 3, Inches(0.5), Inches(2.68), Inches(9.0), Inches(2.55)).table
    tbl.columns[0].width = Inches(4.2)
    tbl.columns[1].width = Inches(1.6)
    tbl.columns[2].width = Inches(3.2)
    for ri, (c0, c1, c2, is_hdr) in enumerate(rows_data):
        for ci, text in enumerate([(c0,c1,c2)[i] for i in range(3)]):
            cell = tbl.cell(ri, ci)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Calibri"
            run.font.size = Pt(11)
            run.font.bold = is_hdr or ci == 2
            if is_hdr:
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                run.font.color.rgb = NAVY if ci < 2 else GOLD
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if ri % 2 == 1 else WHITE

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 4 — WHY SELL NOW
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Market Timing", "Why Now Is the Right Time to Sell")

    add_rect(s, 0.5, 1.42, 9.0, 0.7, LTGOLD, GOLD, Pt(1.5))
    add_textbox(s, 0.65, 1.5, 8.7, 0.54,
                "Q1 seasonal dip = optimal listing window before Q2/Q3 demand surge returns",
                font_size=12.5, bold=True, italic=True, color=NAVY)

    add_bullets(s, [
        "SORA is declining → buyer affordability improving → larger qualified buyer pool",
        "Singapore ultra-wealthy (US$30M+) grew 55% in 5 years: 4,642 → 7,171 in 2026",
        "Knight Frank: UHNWI count forecast to grow a further 46% to 10,497 by 2031",
    ], 0.5, 2.22, 4.5, 2.95)

    add_bullets(s, [
        "25,000–30,000 new citizenships granted annually → steady UHNWI demand pipeline",
        "60% ABSD for foreigners keeps landed as local/PR exclusive — buyers compete in a restricted, motivated pool",
        "2025 landed transaction value up 15% YoY — momentum is building into 2026",
    ], 5.0, 2.22, 4.5, 2.95)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 5 — WHO IS BUYING
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Buyer Profiling", "Who Is Buying Landed in 2026?")

    add_rect(s, 0.5, 1.42, 3.0, 3.8, NAVY)
    add_textbox(s, 0.5, 1.85, 3.0, 0.95, "57.9%",
                font_size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.55, 2.85, 2.9, 0.58,
                "of all landed transactions\npriced above $5M",
                font_size=11.5, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.55, 3.55, 2.9, 0.58,
                "Only 2.1% are HDB upgraders\n(affordability barrier)",
                font_size=10.5, italic=True, color=LTBLUE, align=PP_ALIGN.CENTER)

    segs = [
        ("Multi-generational family wealth",  "Upgrading within the landed tier; wealth preservation mindset."),
        ("UHNWI / New Citizens",              "High-net-worth PRs & new citizens; eligible pool +46% by 2031."),
        ("Condo upgraders",                   "2026 new launches skew 1–2 bed; landlocked condo owners look to landed."),
        ("Developer / redevelopment buyers",  "Buying to rebuild; broadens your buyer pool and exit options."),
    ]
    for i, (title, desc) in enumerate(segs):
        y = 1.42 + i * 0.96
        add_rect(s, 3.7, y, 5.8, 0.85, LIGHT if i%2==0 else LTGOLD, RULE, Pt(0.5))
        add_textbox(s, 3.85, y+0.07, 5.5, 0.3, title, font_size=11.5, bold=True, color=NAVY)
        add_textbox(s, 3.85, y+0.4,  5.5, 0.36, desc,  font_size=10,  color=MIDGRAY)

    add_textbox(s, 3.7, 5.27, 5.8, 0.25,
                "Key insight: the buyer pool is smaller but wealthier and more determined than ever.",
                font_size=9.5, italic=True, color=MIDGRAY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 6 — VALUE DRIVERS
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Valuation", "What Drives Your Home's Value")

    drivers = [
        ("01", "Condition of Home",          "Age, renovation quality, quality of fittings and finishes"),
        ("02", "Facing & Views",             "Landmark views, park/greenery frontage, N-S orientation preferred"),
        ("03", "Parking Availability",       "Number of car porch spaces; covered vs uncovered"),
        ("04", "MRT & Connectivity",         "Walking distance to MRT, bus interchange, expressway access"),
        ("05", "School Catchment",           "Primary school within 1km (P1 registration priority)"),
        ("06", "Comparable PSF Transactions","Recent transacted prices on same street and estate"),
        ("07", "Neighbourhood Upgrading",    "HIP, road widening, new amenities — signals area trajectory"),
    ]
    for i, (num, title, desc) in enumerate(drivers):
        col = 0 if i < 4 else 1
        row = i if i < 4 else i - 4
        x = 0.5 if col == 0 else 5.1
        y = 1.42 + row * 1.02
        bg = LIGHT if col == 0 else LTGOLD
        add_rect(s, x, y, 4.45, 0.88, bg, RULE, Pt(0.5))
        add_textbox(s, x+0.12, y+0.07, 0.5, 0.36, num, font_size=15, bold=True, color=GOLD)
        add_textbox(s, x+0.65, y+0.07, 3.68, 0.32, title, font_size=11.5, bold=True, color=NAVY)
        add_textbox(s, x+0.65, y+0.43, 3.68, 0.36, desc,  font_size=10,   color=MIDGRAY)

    # Bonus card
    bx, by = 5.1, 1.42 + 3 * 1.02
    add_rect(s, bx, by, 4.45, 0.88, NAVY)
    add_textbox(s, bx+0.12, by+0.08, 1.0, 0.28, "+ BONUS", font_size=9, bold=True, color=GOLD)
    add_textbox(s, bx+0.12, by+0.38, 4.2, 0.42,
                "Road Line Category & Redevelopment Potential",
                font_size=11, color=WHITE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 7 — PSF BENCHMARKS
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Pricing Intelligence", "Current PSF Benchmarks by Zone (2026)")

    zones = [
        ("OCR",  "Terrace (Intermediate)",  "$1,100 – $1,400",   "psf land"),
        ("OCR",  "Semi-Detached",           "$1,200 – $1,600",   "psf land"),
        ("RCR",  "Terrace / Semi-D",        "$1,500 – $2,000",   "psf land"),
        ("CCR",  "Detached / GCB",          "$2,000 – $3,500+",  "psf land"),
    ]
    for i, (zone, ptype, psf, note) in enumerate(zones):
        x = 0.5 + i * 2.35
        dark = (i == 3)
        bg = NAVY if dark else LIGHT
        add_rect(s, x, 1.42, 2.15, 2.6, bg, GOLD if dark else RULE, Pt(1 if dark else 0.5))
        add_textbox(s, x+0.08, 1.52, 1.98, 0.42, zone,
                    font_size=22, bold=True, color=GOLD if dark else NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.08, 1.96, 1.98, 0.42, ptype,
                    font_size=10.5, color=LTBLUE if dark else MIDGRAY, align=PP_ALIGN.CENTER)
        add_rect(s, x+0.18, 2.44, 1.78, 0.04, GOLD)
        add_textbox(s, x+0.08, 2.52, 1.98, 0.5, psf,
                    font_size=13.5, bold=True, color=WHITE if dark else NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.08, 3.04, 1.98, 0.28, note,
                    font_size=9, color=LTBLUE if dark else MIDGRAY, align=PP_ALIGN.CENTER)

    add_rect(s, 0.5, 4.12, 9.0, 0.82, LTGOLD, GOLD, Pt(1))
    add_textbox(s, 0.65, 4.2, 8.7, 0.65,
                "Freehold premium: +10–20% over 99-yr leasehold   |   "
                "Corner terrace premium: +15–25% over intermediate   |   "
                "Redevelopment potential adds 5–15% to buyer appetite",
                font_size=10.5, color=NAVY)
    add_textbox(s, 0.5, 5.02, 9.0, 0.3,
                "[INSERT: Your property's zone and indicative PSF range — to be completed before consultation]",
                font_size=9, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 8 — NET PROCEEDS
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Financial Planning", "Your Estimated Net Proceeds")

    rows_8 = [
        ("Estimated Sale Price",                                 "$[AMOUNT]",   False, False),
        ("Less: Agent Commission (2% + 9% GST)",                 "−$[AMOUNT]",  False, True),
        ("Less: CPF Refund + Accrued Interest (if applicable)",  "−$[AMOUNT]",  False, True),
        ("Less: Outstanding Mortgage (if any)",                  "−$[AMOUNT]",  False, True),
        ("Less: Seller's Stamp Duty (SSD — if held <3 years)",   "−$[AMOUNT]",  False, True),
        ("Estimated Cash Proceeds",                              "$[AMOUNT]",   True,  False),
    ]
    tbl2 = s.shapes.add_table(6, 2, Inches(1.0), Inches(1.45), Inches(8.0), Inches(3.5)).table
    tbl2.columns[0].width = Inches(5.8)
    tbl2.columns[1].width = Inches(2.2)
    for ri, (label, val, total, deduct) in enumerate(rows_8):
        for ci, text in enumerate([label, val]):
            cell = tbl2.cell(ri, ci)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Calibri"
            run.font.size = Pt(12)
            run.font.bold = total
            if total:
                run.font.color.rgb = WHITE if ci == 0 else GOLD
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif deduct:
                run.font.color.rgb = NAVY if ci == 0 else RGBColor(0xC0,0x39,0x2B)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            else:
                run.font.color.rgb = NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = LTGOLD

    add_rect(s, 1.0, 5.02, 8.0, 0.38, PLHOLD, RULE, Pt(0.5))
    add_textbox(s, 1.1, 5.06, 7.8, 0.28,
                "Note: BSD (Buyer's Stamp Duty) is paid by the buyer, not the seller.  "
                "SSD: Tier 1 (<1yr) 12% | Tier 2 (1–2yr) 8% | Tier 3 (2–3yr) 4%.",
                font_size=9, color=MIDGRAY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 9 — SELLING JOURNEY
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Process", "The Selling Journey — What to Expect")

    # Timeline bar
    add_rect(s, 0.5, 2.2, 9.0, 0.06, GOLD)

    stages = [
        ("01", "Pre-Launch",       "2–4 weeks",   ["Staging & property prep", "Professional photography", "Drone / aerial video", "Pricing strategy"]),
        ("02", "Active Marketing", "4–8 weeks",   ["PropertyGuru listing", "ERA network activation", "Buyer outreach", "Viewings & feedback"]),
        ("03", "OTP & Negotiation","1–2 weeks",   ["Offer received", "Option to Purchase issued", "14-day exercise period", "Terms negotiated"]),
        ("04", "Completion",       "8–12 weeks",  ["Conveyancing by lawyers", "CPF refund processed", "Loan redemption", "Key handover"]),
    ]
    for i, (num, title, weeks, bullets) in enumerate(stages):
        x = 0.5 + i * 2.35
        # dot
        add_oval(s, x+0.78, 2.01, 0.42, 0.42, NAVY, GOLD)
        add_textbox(s, x+0.78, 2.01, 0.42, 0.42, num,
                    font_size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        # card
        bg = LIGHT if i % 2 == 0 else LTGOLD
        add_rect(s, x, 2.6, 2.15, 2.72, bg, RULE, Pt(0.5))
        add_textbox(s, x+0.1, 2.68, 1.95, 0.34, title, font_size=12.5, bold=True, color=NAVY)
        add_textbox(s, x+0.1, 3.02, 1.95, 0.28, weeks, font_size=9.5, italic=True, color=GOLD)
        add_bullets(s, bullets, x+0.1, 3.35, 1.95, 1.85, font_size=10)

    add_textbox(s, 0.5, 5.33, 9.0, 0.24,
                "Typical total timeline: 3–6 months from signed listing to key handover.",
                font_size=9.5, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 10 — MARKETING
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Marketing Strategy", "How I Will Market Your Home")

    add_rect(s, 0.5,  1.42, 4.45, 3.9, NAVY)
    add_rect(s, 5.05, 1.42, 4.45, 3.9, LIGHT, RULE, Pt(0.5))

    add_textbox(s, 0.65, 1.52, 4.1, 0.5, "PHOTOGRAPHY & MEDIA",
                font_size=12.5, bold=True, color=GOLD)
    add_bullets(s, [
        "360° interior photography with professional gimbal rig",
        "Architectural highlights: staircase, marble, pool, lifestyle spaces",
        "Drone / aerial video showcasing land area and surroundings",
        "Professional video walkthrough for online listings",
        "Staged photography — room-by-room with natural lighting",
    ], 0.65, 2.1, 4.1, 3.0, font_size=11)
    # override bullet colour for dark background manually
    # (bullets in dark box — we'll use white via separate method)
    # Add white text bullets manually
    txBox = s.shapes[-1]
    for p in txBox.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = RGBColor(0xE0, 0xE8, 0xF5)

    add_textbox(s, 5.2, 1.52, 4.1, 0.5, "DISTRIBUTION & OUTREACH",
                font_size=12.5, bold=True, color=NAVY)
    add_bullets(s, [
        "PropertyGuru Premium + ERA PropTrack listing",
        "ERA Singapore network: 7,000+ active agents",
        "ERA Asia-Pacific network for UHNWI & overseas PR buyers",
        "Targeted door-knock outreach to direct buyers in your estate",
        "Instagram & Facebook social media campaign",
        "Weekly seller report: viewings, enquiry quality, market feedback",
    ], 5.2, 2.1, 4.1, 3.0, font_size=11)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 11 — TECHNICAL ADVISORY
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Specialist Knowledge", "Technical Advisory — Why a Specialist Matters")

    add_rect(s, 0.5, 1.42, 9.0, 0.62, LTGOLD, GOLD, Pt(1.2))
    add_textbox(s, 0.65, 1.5, 8.7, 0.48,
                "Serious buyers — especially redevelopment buyers — ask technical questions. "
                "An agent who answers them on the spot wins the listing and closes faster.",
                font_size=10.5, italic=True, color=NAVY)

    tech_items = [
        [
            ("Road Line Category (CAT 1–5)",  "Determines required front setback; CAT 1–2 roads may require land take-back, affecting buildable area."),
            ("SIP / DIP / WSP",               "Sewerage Information Plan, Drainage Interpretation Plan, Water Service Plan — critical pre-redevelopment checks."),
            ("Tree Conservation Areas (TCA)", "South/Eastern zones; removal costs $150–$600/tree; must be considered in site planning."),
        ],
        [
            ("Redevelopment Cost Guide",   "Rebuild from scratch: $800K+, 1.5–2.5 years. Reconstruction: 1–2 years. A&A: 6 months+."),
            ("Pool & Lift Costs",          "Concrete pool: $150K–$300K | Fibreglass: $30K–$60K. Home lift installation: $30K–$100K+."),
            ("Envelope Control & Zoning",  "2-storey vs 3-storey rebuild potential; gross plot ratio; URA planning parameters for your property."),
        ],
    ]
    for col, col_items in enumerate(tech_items):
        for row, (title, desc) in enumerate(col_items):
            x = 0.5 if col == 0 else 5.1
            y = 2.18 + row * 1.08
            add_rect(s, x, y, 4.45, 0.95, LIGHT if col==0 else LTGOLD, RULE, Pt(0.5))
            add_textbox(s, x+0.15, y+0.09, 4.15, 0.3, title, font_size=11.5, bold=True, color=NAVY)
            add_textbox(s, x+0.15, y+0.44, 4.15, 0.44, desc, font_size=10, color=MIDGRAY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 12 — SWAP / NEXT MOVE
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Wealth Planning", "What Happens After You Sell — Your Next Move")

    swap = [
        ("S", "Savings",          "Unlock home equity tied up in your property. Redeploy as liquid capital or down-payment for your next move."),
        ("W", "Work Progression", "As your income grows, your borrowing capacity expands — upgrade to the next tier of property."),
        ("A", "Appreciation",     "Strategic property swaps have historically generated $400K–$500K in net gains within 5–7 years."),
        ("P", "Principal Paydown","Each mortgage payment builds equity. At sale, those gains compound into your next purchase."),
    ]
    for i, (letter, word, desc) in enumerate(swap):
        x = 0.5 + i * 2.35
        dark = (i % 2 == 0)
        add_rect(s, x, 1.42, 2.15, 2.62, NAVY if dark else LTGOLD,
                 GOLD, Pt(1))
        add_textbox(s, x+0.08, 1.5,  1.98, 0.88, letter,
                    font_size=50, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.08, 2.4,  1.98, 0.36, word,
                    font_size=12.5, bold=True, color=WHITE if dark else NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.1,  2.82, 1.95, 1.12, desc,
                    font_size=9.5, color=LTBLUE if dark else MIDGRAY)

    opts = [
        ("RIGHT-SIZE", "Sell your landed, move to a condo — unlock $2M+ cash, reduce maintenance burden."),
        ("UPGRADE",    "Proceeds from your landed can fund a GCB or prime district (D10/D11) detached."),
        ("INVEST",     "Deploy proceeds into income-generating properties for wealth accumulation and retirement."),
    ]
    for i, (label, desc) in enumerate(opts):
        x = 0.5 + i * 3.1
        add_rect(s, x, 4.18, 2.9, 1.1, LIGHT, RULE, Pt(0.5))
        add_textbox(s, x+0.12, 4.26, 2.65, 0.3, label,
                    font_size=10.5, bold=True, color=GOLD)
        add_textbox(s, x+0.12, 4.6,  2.65, 0.58, desc,
                    font_size=9.5, color=NAVY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 13 — TRACK RECORD
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "My Credentials", "My Track Record in Landed Property")

    add_textbox(s, 0.5, 1.45, 5.5, 0.28, "RECENT TRANSACTIONS",
                font_size=9.5, bold=True, color=GOLD)

    tbl3 = s.shapes.add_table(4, 4, Inches(0.5), Inches(1.8), Inches(5.5), Inches(1.95)).table
    tbl3.columns[0].width = Inches(2.1)
    tbl3.columns[1].width = Inches(1.3)
    tbl3.columns[2].width = Inches(1.2)
    tbl3.columns[3].width = Inches(0.9)
    hdr_texts = ["Address / District", "Type", "Price", "DOM"]
    row_data  = [
        ["[INSERT: Address]", "[Type]", "$[X]M", "[X] days"],
        ["[INSERT: Address]", "[Type]", "$[X]M", "[X] days"],
        ["[INSERT: Address]", "[Type]", "$[X]M", "[X] days"],
    ]
    for ci, ht in enumerate(hdr_texts):
        cell = tbl3.cell(0, ci)
        cell.text = ht
        p = cell.text_frame.paragraphs[0]
        r = p.runs[0] if p.runs else p.add_run()
        r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for ri, row in enumerate(row_data):
        for ci, text in enumerate(row):
            cell = tbl3.cell(ri+1, ci)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            r = p.runs[0] if p.runs else p.add_run()
            r.font.name = "Calibri"; r.font.size = Pt(11)
            r.font.color.rgb = NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri%2==0 else WHITE

    add_textbox(s, 0.5, 3.88, 5.5, 0.28, "CREDENTIALS",
                font_size=9.5, bold=True, color=GOLD)
    add_bullets(s, [
        "CEA Registered Salesperson (RES)",
        "ERA Landed Expert Series — all 5 modules completed",
        "ERA PropNett Network member",
        "Districts served: [INSERT districts]",
    ], 0.5, 4.2, 5.4, 1.1, font_size=11)

    add_textbox(s, 6.2, 1.45, 3.3, 0.28, "CLIENT TESTIMONIALS",
                font_size=9.5, bold=True, color=GOLD)
    for i in range(2):
        y = 1.8 + i * 1.78
        add_rect(s, 6.2, y, 3.3, 1.58, PLHOLD, RULE, Pt(0.5))
        add_textbox(s, 6.35, y+0.12, 3.0, 1.0,
                    '"[INSERT: Client testimonial — 1–2 sentences about the experience of selling through you.]"',
                    font_size=10, italic=True, color=NAVY)
        add_textbox(s, 6.35, y+1.2, 3.0, 0.28,
                    "— [Client Name], [Type], [District]",
                    font_size=9.5, bold=True, color=MIDGRAY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 14 — 7-POINT PROMISE
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "My Commitment", "My 7-Point Marketing Promise to You")

    promises = [
        ("01", "Photography & Drone Video",        "Professional shoot within 5 days. 360°, gimbal interior, drone aerial, lifestyle video."),
        ("02", "Listing Goes Live in 7 Days",       "PropertyGuru Premium + ERA PropTrack activated within 1 week of signing."),
        ("03", "10+ Qualified Buyer Contacts/Week", "Minimum 10 targeted outreach contacts per week during active campaign."),
        ("04", "Weekly Seller Update",              "Viewings count, enquiry quality, market feedback, next steps — in writing every week."),
        ("05", "Technical Advisory On-Demand",      "Road line, SIP/DIP, redevelopment Q&A — answered on the spot for serious buyers."),
        ("06", "Full Negotiation Representation",   "I represent your interests at OTP stage: pricing, conditions, timeline, deposit."),
        ("07", "After-Sale Support",                "CPF refund coordination, key handover checklist, referrals for solicitors and movers."),
    ]
    for i, (num, title, desc) in enumerate(promises):
        col = 0 if i < 4 else 1
        row = i if i < 4 else i - 4
        x   = 0.5 if col == 0 else 5.1
        y   = 1.42 + row * 1.02
        add_rect(s, x, y, 4.45, 0.88, LIGHT if col==0 else LTGOLD, RULE, Pt(0.5))
        add_oval(s, x+0.12, y+0.22, 0.42, 0.42, NAVY)
        add_textbox(s, x+0.12, y+0.22, 0.42, 0.42, num,
                    font_size=9.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.65, y+0.07, 3.7, 0.3, title, font_size=11.5, bold=True, color=NAVY)
        add_textbox(s, x+0.65, y+0.44, 3.7, 0.36, desc,  font_size=9.5,  color=MIDGRAY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 15 — FAQ
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, WHITE)
    add_label_title(s, "Common Questions", "Frequently Asked Questions")

    faqs = [
        ("Do I need to move out before the property is marketed?",
         "No — we can market with you in residence. Professional staging and photography work around your schedule. Occupied properties often photograph better with furniture."),
        ("How long will it realistically take to sell?",
         "Well-priced landed homes in active districts typically receive offers within 4–8 weeks. Total timeline: 3–6 months including legal completion."),
        ("What if a buyer wants to redevelop my property?",
         "This broadens your buyer pool — redevelopment buyers are often cash-rich and motivated. I advise on redevelopment potential, BCA process, and pricing strategy."),
        ("What happens to my CPF monies when I sell?",
         "CPF OA principal used + accrued interest (at CPF OA rate, currently 2.5% p.a.) must be fully refunded to your CPF OA upon sale. I will walk you through the exact calculation."),
    ]
    for i, (q, a) in enumerate(faqs):
        col = 0 if i < 2 else 1
        row = i if i < 2 else i - 2
        x = 0.5 if col == 0 else 5.1
        y = 1.42 + row * 2.0
        add_rect(s, x, y, 4.45, 1.88, LIGHT if col==0 else LTGOLD, RULE, Pt(0.5))
        add_oval(s, x+0.12, y+0.16, 0.38, 0.38, NAVY)
        add_textbox(s, x+0.12, y+0.16, 0.38, 0.38, "Q",
                    font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.62, y+0.1, 3.72, 0.54, q, font_size=11, bold=True, color=NAVY)
        add_oval(s, x+0.12, y+0.78, 0.38, 0.38, GOLD)
        add_textbox(s, x+0.12, y+0.78, 0.38, 0.38, "A",
                    font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.62, y+0.72, 3.72, 1.08, a, font_size=10, color=MIDGRAY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 16 — NEXT STEPS / BACK COVER
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, NAVY)

    add_textbox(s, 0.5, 0.28, 9.0, 0.65, "Let's Get Started",
                font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 0.96, 9.0, 0.32,
                "Three simple steps to finding out what your home is worth today.",
                font_size=12.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

    steps = [
        ("1", "Free Property Assessment",   "I visit your home, assess condition, run recent comps, and give you an honest indicative valuation. No strings attached."),
        ("2", "Marketing Strategy Meeting", "We agree on the right price, launch timing, and marketing approach. I share the comparable transactions most relevant to your home."),
        ("3", "Sign & Launch",              "We sign the exclusive listing. Photography within 5 days, listing live within 7 days, buyer outreach begins immediately."),
    ]
    for i, (n, title, desc) in enumerate(steps):
        x = 0.5 + i * 3.1
        add_rect(s, x, 1.45, 2.9, 2.65, DARK2, GOLD, Pt(1.2))
        add_oval(s, x+1.15, 1.56, 0.6, 0.6, GOLD)
        add_textbox(s, x+1.15, 1.56, 0.6, 0.6, n,
                    font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.12, 2.28, 2.65, 0.44, title,
                    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x+0.15, 2.78, 2.6, 1.22, desc, font_size=10, color=LTBLUE)

    # Contact bar
    add_rect(s, 0.5, 4.25, 9.0, 1.1, DARK3, RGBColor(0x2E,0x41,0x66), Pt(0.5))
    add_textbox(s, 0.75, 4.36, 3.2, 0.4,  "[INSERT: Agent Name]",        font_size=14.5, bold=True, color=WHITE)
    add_textbox(s, 0.75, 4.78, 3.2, 0.3,  "[INSERT: CEA Reg. No.] | ERA Singapore", font_size=10, color=LTBLUE)
    add_textbox(s, 4.1,  4.36, 2.6, 0.32, "[INSERT: Mobile / WhatsApp]", font_size=12, color=GOLD)
    add_textbox(s, 4.1,  4.72, 2.6, 0.32, "[INSERT: Email Address]",     font_size=12, color=GOLD)
    add_rect(s, 7.6, 4.3, 1.6, 1.0, DARK2, RGBColor(0x2E,0x41,0x66), Pt(0.5))
    add_textbox(s, 7.6, 4.55, 1.6, 0.55, "ERA\nSINGAPORE",
                font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ─── SAVE ─────────────────────────────────────────────────────────────────
    out = "/Users/riopang/Desktop/Claude Code/Landed Seller Pitch Deck.pptx"
    prs.save(out)
    print(f"Saved → {out}")

build()
